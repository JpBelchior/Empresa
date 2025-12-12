import sys
import base64
from io import BytesIO
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement


def add_base64_image(slide, base64_str, left, top, width, height):
    try:
        if ',' in base64_str:
            base64_str = base64_str.split(',')[1]

        img_data = base64.b64decode(base64_str)
        stream = BytesIO(img_data)
        slide.shapes.add_picture(
            stream, Inches(left), Inches(top),
            width=Inches(width), height=Inches(height)
        )
        return True
    except Exception as e:
        print(f"Erro imagem: {e}", file=sys.stderr)
        return False


def create_diagonal_strip(slide, pres):
    """Cria faixa azul diagonal inclinada para esquerda"""
    width = pres.slide_width
    height = pres.slide_height
    
    # Criar shape
    shapes = slide.shapes
    shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(1), Inches(1))
    
    # Ajustar para polígono customizado via XML
    sp = shape._element
    spPr = sp.spPr
    
    # Remover geometria padrão
    for child in list(spPr):
        if 'prstGeom' in child.tag:
            spPr.remove(child)
    
    # Criar geometria customizada
    custGeom = OxmlElement('a:custGeom')
    
    avLst = OxmlElement('a:avLst')
    custGeom.append(avLst)
    
    gdLst = OxmlElement('a:gdLst')
    custGeom.append(gdLst)
    
    ahLst = OxmlElement('a:ahLst')
    custGeom.append(ahLst)
    
    cxnLst = OxmlElement('a:cxnLst')
    custGeom.append(cxnLst)
    
    # Criar path com 4 vértices (paralelogramo inclinado)
    pathLst = OxmlElement('a:pathLst')
    path = OxmlElement('a:path')
    path.set('w', str(width))
    path.set('h', str(height))
    
    # Vértices do paralelogramo (faixa inclinada para esquerda)
    # Vai do topo (mais à esquerda) até o final da página à direita
    top_left_x = int(Inches(6.5))     # Topo começa mais à esquerda
    bottom_left_x = int(Inches(7.5))  # Base começa mais à direita (inclinação)
    
    # Ponto 1: Topo esquerdo
    moveTo = OxmlElement('a:moveTo')
    pt = OxmlElement('a:pt')
    pt.set('x', str(top_left_x))
    pt.set('y', '0')
    moveTo.append(pt)
    path.append(moveTo)
    
    # Ponto 2: Topo direito (final da página)
    lineTo1 = OxmlElement('a:lnTo')
    pt1 = OxmlElement('a:pt')
    pt1.set('x', str(int(width)))  # Vai até o final
    pt1.set('y', '0')
    lineTo1.append(pt1)
    path.append(lineTo1)
    
    # Ponto 3: Base direita (final da página)
    lineTo2 = OxmlElement('a:lnTo')
    pt2 = OxmlElement('a:pt')
    pt2.set('x', str(int(width)))  # Vai até o final
    pt2.set('y', str(int(height)))
    lineTo2.append(pt2)
    path.append(lineTo2)
    
    # Ponto 4: Base esquerda
    lineTo3 = OxmlElement('a:lnTo')
    pt3 = OxmlElement('a:pt')
    pt3.set('x', str(bottom_left_x))
    pt3.set('y', str(int(height)))
    lineTo3.append(pt3)
    path.append(lineTo3)
    
    close = OxmlElement('a:close')
    path.append(close)
    
    pathLst.append(path)
    custGeom.append(pathLst)
    spPr.append(custGeom)
    
    # Definir posição e tamanho
    shape.left = 0
    shape.top = 0
    shape.width = width
    shape.height = height
    
    # === GRADIENTE AZUL (escuro embaixo, claro em cima) === #
    fill = shape.fill
    fill.gradient()
    fill.gradient_angle = 90
    stop_dark = fill.gradient_stops[0]
    stop_dark.position = 0.0
    stop_dark.color.rgb = RGBColor(10, 50, 100)
    stop_light = fill.gradient_stops[1]
    stop_light.position = 1.0
    stop_light.color.rgb = RGBColor(30, 115, 190)
    
    shape.line.fill.background()
    
    return shape


def gerar_capa(pres, dados):
    print("Criando capa corporativa...", file=sys.stderr)
    slide = pres.slides.add_slide(pres.slide_layouts[6])

    # === FUNDO COM GRADIENTE PROFISSIONAL === #
    fill = slide.background.fill
    fill.gradient()
    grad = fill.gradient_stops

    grad[0].position = 0.0
    grad[0].color.rgb = RGBColor(10, 38, 85)

    if len(grad) < 2:
        grad.add_stop(1.0)

    grad[1].position = 1.0
    grad[1].color.rgb = RGBColor(4, 20, 48)

    # === FAIXA LATERAL DIAGONAL (inclinada para esquerda) === #
    create_diagonal_strip(slide, pres)

    # === LOGO DA EMPRESA === #
    logo = dados.get("imagens", {}).get("logo_empresa")
    if logo:
        add_base64_image(slide, logo, 0.35, 0.35, 1.5, 0.92)

    # === TÍTULO PRINCIPAL (movido mais para esquerda) === #
    title_left = Inches(2.0)  # Era 2.8, agora 2.0 (mais à esquerda)
    title_top = Inches(1.5)
    title_w = Inches(6.5)
    title_h = Inches(1.8)

    title_box = slide.shapes.add_textbox(title_left, title_top, title_w, title_h)
    tf = title_box.text_frame
    tf.clear()

    # Primeira linha: ANÁLISE DE
    p1 = tf.paragraphs[0]
    p1.text = "ANÁLISE DE"
    p1.font.name = "Calibri"
    p1.font.size = Pt(50)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(255, 255, 255)

    # Segunda linha: RISCOS
    p2 = tf.add_paragraph()
    p2.text = "RISCOS"
    p2.font.name = "Calibri"
    p2.font.size = Pt(50)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(255, 255, 255)

    # === LINHA DECORATIVA DUPLA (mantidas retas) === #
    # Linha superior (fina)
    linha1 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        title_left,
        title_top + Inches(1.8),
        Inches(4.0),
        Inches(0.02)
    )
    linha1.fill.solid()
    linha1.fill.fore_color.rgb = RGBColor(255, 255, 255)
    linha1.fill.transparency = 0.3
    linha1.line.fill.background()

    # Linha inferior (grossa)
    linha2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        title_left,
        title_top + Inches(1.88),
        Inches(4.0),
        Inches(0.05)
    )
    linha2.fill.solid()
    linha2.fill.fore_color.rgb = RGBColor(255, 255, 255)
    linha2.line.fill.background()

    # === INFORMAÇÕES DA EMPRESA === #
    empresa = dados.get("dados", {}).get("nome_empresa", "")
    local = dados.get("dados", {}).get("localizacao_analise", "")
    
    info_top = title_top + Inches(2.05)
    
    # Nome da empresa
    if empresa:
        empresa_box = slide.shapes.add_textbox(
            title_left,
            info_top,
            Inches(6.0),
            Inches(0.4)
        )
        tf_emp = empresa_box.text_frame
        tf_emp.clear()
        p_emp = tf_emp.paragraphs[0]
        p_emp.text = empresa.upper()
        p_emp.font.name = "Calibri"
        p_emp.font.size = Pt(22)
        p_emp.font.bold = True
        p_emp.font.color.rgb = RGBColor(255, 255, 255)

    # Localização
    if local:
        loc_box = slide.shapes.add_textbox(
            title_left,
            info_top + Inches(0.45),
            Inches(6.0),
            Inches(0.3)
        )
        tf_loc = loc_box.text_frame
        tf_loc.clear()
        p_loc = tf_loc.paragraphs[0]
        p_loc.text = local.upper()
        p_loc.font.name = "Calibri Light"
        p_loc.font.size = Pt(16)
        p_loc.font.color.rgb = RGBColor(200, 200, 200)

    # === DATA (MÊS / ANO em duas linhas) === #
    agora = datetime.now()
    
   
    meses_pt = {
        'January': 'JANEIRO',
        'February': 'FEVEREIRO',
        'March': 'MARÇO',
        'April': 'ABRIL',
        'May': 'MAIO',
        'June': 'JUNHO',
        'July': 'JULHO',
        'August': 'AGOSTO',
        'September': 'SETEMBRO',
        'October': 'OUTUBRO',
        'November': 'NOVEMBRO',
        'December': 'DEZEMBRO'
    }
    
    mes_ingles = agora.strftime("%B")
    mes = meses_pt.get(mes_ingles, mes_ingles.upper())  # Mês em português
    ano = agora.strftime("%Y")          # 2025
    ano_espacado = " ".join(ano)        
    
    data_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.7),  
        Inches(2.0), Inches(0.6)
    )
    tf_d = data_box.text_frame
    tf_d.clear()
    
    # Primeira linha: MÊS (em português)
    p_mes = tf_d.paragraphs[0]
    p_mes.text = mes
    p_mes.font.name = "Calibri"
    p_mes.font.size = Pt(16)
    p_mes.font.bold = True
    p_mes.font.color.rgb = RGBColor(255, 255, 255)
    
    # Segunda linha: ANO (com espaçamento entre dígitos)
    p_ano = tf_d.add_paragraph()
    p_ano.text = ano_espacado
    p_ano.font.name = "Arial"
    p_ano.font.size = Pt(24)
    p_ano.font.bold = True
    p_ano.font.color.rgb = RGBColor(255, 255, 255)

    print("Capa corporativa criada!", file=sys.stderr)
    return slide