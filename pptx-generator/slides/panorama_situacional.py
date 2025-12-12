#!/usr/bin/env python3
# -*- coding: utf-8 -*-
import sys
import os
import base64
from io import BytesIO
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


def add_base64_image(slide, base64_str, left, top, width, height):
    """Adiciona imagem base64 ao slide"""
    try:
        if base64_str and base64_str.startswith('data:image'):
            base64_data = base64_str.split(',')[1]
        else:
            base64_data = base64_str
        
        image_data = base64.b64decode(base64_data)
        image_stream = BytesIO(image_data)
        
        slide.shapes.add_picture(
            image_stream,
            Inches(left),
            Inches(top),
            width=Inches(width),
            height=Inches(height)
        )
        return True
    except Exception as e:
        print(f"Erro ao adicionar imagem base64: {e}", file=sys.stderr)
        return False


def add_local_image(slide, image_path, left, top, width, height):
    """Adiciona imagem local ao slide"""
    try:
        if os.path.exists(image_path):
            slide.shapes.add_picture(
                image_path, 
                Inches(left), 
                Inches(top),
                width=Inches(width), 
                height=Inches(height)
            )
            return True
        else:
            print(f"Imagem não encontrada: {image_path}", file=sys.stderr, flush=True)
            return False
    except Exception as e:
        print(f"Erro ao adicionar imagem local: {e}", file=sys.stderr, flush=True)
        return False


def criar_cabecalho_slide(slide, pres, numero, titulo):
    """Cria o cabeçalho padrão com paralelogramo azul e cinza"""
    
    # === PARALELOGRAMO AZUL (número) === #
    para_azul = slide.shapes.add_shape(
        MSO_SHAPE.PARALLELOGRAM,
        Inches(0.2),    # Posição à esquerda
        Inches(0.45),   # Descido 0.2 inches
        Inches(0.5),    # Largura menor
        Inches(0.5)     # Altura menor
    )
    
    para_azul.fill.solid()
    para_azul.fill.fore_color.rgb = RGBColor(0, 102, 204)
    para_azul.line.fill.background()
    para_azul.rotation = 0
    
    # Adicionar número
    tf = para_azul.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = str(numero)
    p.font.name = "Arial"
    p.font.size = Pt(20)  # Reduzido de 32 para 20
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # === PARALELOGRAMO CINZA (título) === #
    para_cinza = slide.shapes.add_shape(
        MSO_SHAPE.PARALLELOGRAM,
        Inches(0.85),   # Logo após o paralelogramo azul
        Inches(0.45),   # Alinhado verticalmente
        Inches(3.5),    # Largura para o texto (um pouco maior para "PANORAMA SITUACIONAL")
        Inches(0.5)     # Altura menor
    )
    
    para_cinza.fill.solid()
    para_cinza.fill.fore_color.rgb = RGBColor(64, 64, 64)
    para_cinza.line.fill.background()
    para_cinza.rotation = 0
    
    # Adicionar título
    tf = para_cinza.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = titulo
    p.font.name = "Arial"
    p.font.size = Pt(14)  # Reduzido de 20 para 14
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.LEFT
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.2)


def gerar_panorama_situacional(pres, dados):
    """Gera o slide de Panorama Situacional"""
    print("Criando slide de Panorama Situacional...", file=sys.stderr)
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    
    # === FUNDO BRANCO === #
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # === CABEÇALHO (Paralelogramo azul + Título) === #
    criar_cabecalho_slide(slide, pres, 3, "PANORAMA SITUACIONAL")
    
    # === DADOS DO RELATÓRIO === #
    dados_info = dados.get("dados", {})
    imagens_info = dados.get("imagens", {})
    
    empresa = dados.get("dados", {}).get("nome_empresa", "")
    nome_local = dados_info.get("localizacao_analise", "NOME DO LOCAL")
    panorama_texto = dados_info.get("panorama", "")
    referencias_proximas = dados_info.get("referencias_proximas", "")
    referencias_lista = dados_info.get("referencias_proximas_lista", [])
    imagem_area = imagens_info.get("imagem_area", "")
    
    # === TÍTULO "NOME DO LOCAL" (centralizado, acima da imagem) === #
    titulo_local_box = slide.shapes.add_textbox(
        Inches(5.2),   # Alinhado com a borda esquerda da imagem
        Inches(1.2),  # Subido
        Inches(4.5),   # Mesma largura da imagem
        Inches(0.4)  # Altura reduzida
    )
    tf = titulo_local_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = empresa.upper()
    p.font.name = "Arial"
    p.font.size = Pt(18)  # Tamanho reduzido de 24 para 18
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER
    
    # === IMAGEM DO LOCAL (metade direita do slide) === #
    if imagem_area:
        # Adicionar imagem centralizada no lado direito
        add_base64_image(
            slide,
            imagem_area,
            left=5.2,   # Metade direita
            top=1.6,    # Subido
            width=4.5,
            height=3.7
        )
    else:
        # Placeholder se não houver imagem
        placeholder = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(5.2),
            Inches(1.4),  # Subido
            Inches(4.5),
            Inches(3.2)
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(200, 200, 200)
        placeholder.line.color.rgb = RGBColor(0, 102, 204)
        placeholder.line.width = Pt(2)
        
        tf = placeholder.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = "IMAGEM DO LOCAL"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(100, 100, 100)
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # === LADO ESQUERDO - DADOS DO PANORAMA === #
    
    # 1. TEXTO PANORAMA SITUACIONAL - EXPOSIÇÃO AO RISCO
    if panorama_texto:
        texto_box = slide.shapes.add_textbox(
            Inches(1.0),  # Movido 0.4" para direita (era 0.6)
            Inches(1.4),  # Subido de 1.8 para 1.4
            Inches(3.5),  # Reduzido (era 3.9)
            Inches(0.7)  # Altura reduzida (era 1.0)
        )
        tf = texto_box.text_frame
        tf.clear()
        tf.word_wrap = True
        
        # Título da seção
        p_titulo = tf.paragraphs[0]
        p_titulo.text = "EXPOSIÇÃO AO RISCO"
        p_titulo.font.name = "Arial"
        p_titulo.font.size = Pt(9)
        p_titulo.font.bold = True
        p_titulo.font.color.rgb = RGBColor(0, 51, 102)
        p_titulo.alignment = PP_ALIGN.LEFT
        
        # Texto do panorama
        p_texto = tf.add_paragraph()
        p_texto.text = panorama_texto
        p_texto.font.name = "Arial"
        p_texto.font.size = Pt(9)
        p_texto.font.color.rgb = RGBColor(80, 80, 80)
        p_texto.alignment = PP_ALIGN.LEFT
        p_texto.space_before = Pt(6)
    
    # 2. OBSERVAÇÕES DA LOCALIDADE / REFERÊNCIAS PRÓXIMAS
    ref_texto_box = slide.shapes.add_textbox(
        Inches(1.0),  # Movido 0.4" para direita (era 0.6)
        Inches(2.3),  # Aproximado do texto anterior (era 2.5)
        Inches(3.5),  # Ajustado (era 4.0)
        Inches(0.8)  # Altura reduzida (era 1.2)
    )
    tf = ref_texto_box.text_frame
    tf.clear()
    tf.word_wrap = True
    
    # ÍCONE DE CRIMINALIDADE (ao lado esquerdo do título)
    current_dir = os.path.dirname(__file__)  # pptx-generator/slides/
    parent_dir = os.path.dirname(current_dir)  # pptx-generator/
    icon_criminalidade_path = os.path.join(parent_dir, "images", "ICON_CRIMINALIDADE.png")
    
    add_local_image(
        slide,
        icon_criminalidade_path,
        left=0.3,
        top=2.3,  # Ajustado para acompanhar o texto (era 2.5)
        width=0.6,
        height=0.6
    )
    
    # Título da seção
    p_titulo = tf.paragraphs[0]
    p_titulo.text = "OBSERVAÇÕES DA LOCALIDADE (DADOS, TAXAS DE CRIMINALIDADE ETC.)"
    p_titulo.font.name = "Arial"
    p_titulo.font.size = Pt(9)
    p_titulo.font.bold = True
    p_titulo.font.color.rgb = RGBColor(0, 51, 102)
    p_titulo.alignment = PP_ALIGN.LEFT
    
    # Se tiver lista de referências
    if referencias_lista and len(referencias_lista) > 0:
        for ref in referencias_lista:
            p = tf.add_paragraph()
            p.text = f"• {ref}"
            p.font.name = "Arial"
            p.font.size = Pt(9)
            p.font.color.rgb = RGBColor(80, 80, 80)
            p.alignment = PP_ALIGN.LEFT
            p.space_before = Pt(3)
    elif referencias_proximas:
        # Se tiver texto de referências
        p = tf.add_paragraph()
        p.text = referencias_proximas
        p.font.name = "Arial"
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(80, 80, 80)
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(6)
    
    # 3. LOCALIZAÇÃO (NOME DO LOCAL + ENDEREÇO)
    loc_texto_box = slide.shapes.add_textbox(
        Inches(1.0),  # Movido 0.4" para direita (era 0.6)
        Inches(3.3),  # Aproximado do texto anterior (era 3.8)
        Inches(3.5),  # Ajustado (era 4.0)
        Inches(0.8)  # Altura reduzida (era 1.0)
    )
    tf = loc_texto_box.text_frame
    tf.clear()
    tf.word_wrap = True
    
    # ÍCONE DE LOCALIZAÇÃO (ao lado esquerdo do título)
    icon_localizacao_path = os.path.join(parent_dir, "images", "ICON_LOCALIZAÇÃO.png")
    
    add_local_image(
        slide,
        icon_localizacao_path,
        left=0.3,
        top=3.3,  # Ajustado para acompanhar o texto (era 3.8)
        width=0.6,
        height=0.6
    )
    
    # Título LOCALIZAÇÃO
    p_loc_titulo = tf.paragraphs[0]
    p_loc_titulo.text = "LOCALIZAÇÃO"
    p_loc_titulo.font.name = "Arial"
    p_loc_titulo.font.size = Pt(9)
    p_loc_titulo.font.bold = True
    p_loc_titulo.font.color.rgb = RGBColor(0, 51, 102)
    p_loc_titulo.alignment = PP_ALIGN.LEFT
    
    # Nome do local
    p_nome = tf.add_paragraph()
    p_nome.text = nome_local
    p_nome.font.name = "Arial"
    p_nome.font.size = Pt(9)
    p_nome.font.bold = True
    p_nome.font.color.rgb = RGBColor(80, 80, 80)
    p_nome.alignment = PP_ALIGN.LEFT
    p_nome.space_before = Pt(6)
    
    # Endereço (placeholder)
    p_endereco = tf.add_paragraph()
    p_endereco.text = "ENDEREÇO"
    p_endereco.font.name = "Arial"
    p_endereco.font.size = Pt(9)
    p_endereco.font.color.rgb = RGBColor(80, 80, 80)
    p_endereco.alignment = PP_ALIGN.LEFT
    p_endereco.space_before = Pt(3)
    
    caixa_azul = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.3),    # Começa do lado esquerdo
    Inches(4.3),    # Abaixo da seção de localização
    Inches(4.5),    # Até um pouco antes da metade da página
    Inches(1.2)     # Altura da caixa
)

# Fundo branco
    caixa_azul.fill.solid()
    caixa_azul.fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Borda azul (0,102,204)
    caixa_azul.line.fill.solid()
    caixa_azul.line.color.rgb = RGBColor(0, 102, 204)
    caixa_azul.line.width = Pt(2) 
    
    print("Slide de Panorama Situacional criado!", file=sys.stderr)
    return slide