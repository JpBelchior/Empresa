#!/usr/bin/env python3
# -*- coding: utf-8 -*-
import sys
import os
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import base64
from io import BytesIO


def add_local_image(slide, image_path, left, top, width, height):
    """Adiciona imagem local ao slide"""
    try:
        if os.path.exists(image_path):
            slide.shapes.add_picture(
                image_path, 
                Inches(left), 
                Inches(top),
                width=Inches(width), 
                height=Inches(height)
            )
            return True
        else:
            print(f"Imagem não encontrada: {image_path}", file=sys.stderr, flush=True)
            return False
    except Exception as e:
        print(f"Erro ao adicionar imagem local: {e}", file=sys.stderr, flush=True)
        return False


def criar_slide_vulnerabilidades(pres, dados, itens_slide):

    slide = pres.slides.add_slide(pres.slide_layouts[6])

    # === FUNDO BRANCO === #
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # === PARALELOGRAMO AZUL === #
    para_azul = slide.shapes.add_shape(
        MSO_SHAPE.PARALLELOGRAM,
        Inches(0.2), Inches(0.2), Inches(0.5), Inches(0.5)
    )
    para_azul.fill.solid()
    para_azul.fill.fore_color.rgb = RGBColor(30, 115, 190)
    para_azul.line.fill.background()

    tf_azul = para_azul.text_frame
    tf_azul.text = "8"
    p = tf_azul.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    tf_azul.vertical_anchor = MSO_ANCHOR.MIDDLE

    # === PARALELOGRAMO CINZA === #
    para_cinza = slide.shapes.add_shape(
        MSO_SHAPE.PARALLELOGRAM,
        Inches(0.85), Inches(0.2), Inches(5.0), Inches(0.5)
    )
    para_cinza.fill.solid()
    para_cinza.fill.fore_color.rgb = RGBColor(70, 70, 70)
    para_cinza.line.fill.background()

    tf_cinza = para_cinza.text_frame
    tf_cinza.text = "VULNERABILIDADES ESPECÍFICAS"
    tf_cinza.margin_left = Inches(0.2)
    tf_cinza.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_cinza.paragraphs[0].font.bold = True
    tf_cinza.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # === SUBTÍTULO === #
    subtitulo = slide.shapes.add_textbox(
        Inches(2.5), Inches(0.75), Inches(5.0), Inches(0.35)
    )
    p = subtitulo.text_frame.paragraphs[0]
    p.text = "Áreas Restritas"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 38, 77)
    p.alignment = PP_ALIGN.CENTER
