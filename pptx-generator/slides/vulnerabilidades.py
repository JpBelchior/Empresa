#!/usr/bin/env python3
# -*- coding: utf-8 -*-
import sys
import os
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import base64
from io import BytesIO


def add_local_image(slide, image_path, left, top, width, height):
    """Adiciona imagem local ao slide"""
    try:
        if os.path.exists(image_path):
            slide.shapes.add_picture(
                image_path, 
                Inches(left), 
                Inches(top),
                width=Inches(width), 
                height=Inches(height)
            )
            return True
        else:
            print(f"Imagem não encontrada: {image_path}", file=sys.stderr, flush=True)
            return False
    except Exception as e:
        print(f"Erro ao adicionar imagem local: {e}", file=sys.stderr, flush=True)
        return False


def criar_slide_vulnerabilidades(pres, dados, itens_slide):

    slide = pres.slides.add_slide(pres.slide_layouts[6])

    # === FUNDO BRANCO === #
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # === PARALELOGRAMO AZUL === #
    para_azul = slide.shapes.add_shape(
        MSO_SHAPE.PARALLELOGRAM,
        Inches(0.2), Inches(0.2), Inches(0.5), Inches(0.5)
    )
    para_azul.fill.solid()
    para_azul.fill.fore_color.rgb = RGBColor(30, 115, 190)
    para_azul.line.fill.background()

    tf_azul = para_azul.text_frame
    tf_azul.text = "7"
    p = tf_azul.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    tf_azul.vertical_anchor = MSO_ANCHOR.MIDDLE

    # === PARALELOGRAMO CINZA === #
    para_cinza = slide.shapes.add_shape(
        MSO_SHAPE.PARALLELOGRAM,
        Inches(0.85), Inches(0.2), Inches(5.0), Inches(0.5)
    )
    para_cinza.fill.solid()
    para_cinza.fill.fore_color.rgb = RGBColor(70, 70, 70)
    para_cinza.line.fill.background()

    tf_cinza = para_cinza.text_frame
    tf_cinza.text = "VULNERABILIDADES"
    tf_cinza.margin_left = Inches(0.2)
    tf_cinza.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_cinza.paragraphs[0].font.bold = True
    tf_cinza.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # === SUBTÍTULO === #
    subtitulo = slide.shapes.add_textbox(
        Inches(2.5), Inches(0.75), Inches(5.0), Inches(0.35)
    )
    p = subtitulo.text_frame.paragraphs[0]
    p.text = "Não Conformidades"
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = RGBColor(0, 38, 77)
    p.alignment = PP_ALIGN.CENTER

    # ================= TABELA ================= #
    rows = 1 + len(itens_slide)
    cols = 7

    table = slide.shapes.add_table(
        rows, cols,
        Inches(0.3), Inches(1.2),
        Inches(9.0),
        Inches(0.2 + 0.6 * len(itens_slide))
    ).table

    # Larguras
    widths = [1.0, 0.5, 2.0, 1.7, 1.2, 1.3, 1.0]
    for i, w in enumerate(widths):
        table.columns[i].width = Inches(w)

    # Alturas
    table.rows[0].height = Inches(0.3)
    for i in range(1, rows):
        table.rows[i].height = Inches(0.6)

    # Cabeçalho
    headers = [
        "Elementos", "NC", "Não Conformidade",
        "Observação", "Risco", "Localização", "Criticidade"
    ]

    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(30, 115, 190)

        p = cell.text_frame.paragraphs[0]
        p.text = header
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    
    for row_idx, item in enumerate(itens_slide, start=1):
        valores = [
            item["elemento"],
            item["nc"],
            item["nao_conformidade"],
            item["observacao"],
            item["risco"],
            item["localizacao"],
            item["criticidade"],
        ]

        for col_idx, valor in enumerate(valores):
            cell = table.cell(row_idx, col_idx)
            p = cell.text_frame.paragraphs[0]
            p.text = valor
            p.font.size = Pt(10)
            p.alignment = PP_ALIGN.CENTER
            cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    # -------- TÍTULO DA LEGENDA  -------- #
    titulo_legenda = slide.shapes.add_textbox(
        Inches(0.3),
        Inches(4.55),   # um pouco acima da caixa
        Inches(4.2),
        Inches(0.25)
    )

    p = titulo_legenda.text_frame.paragraphs[0]
    p.text = "Criticidade Risco"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.LEFT

    legenda = slide.shapes.add_textbox(
        Inches(0.2),        # Left
        Inches(4.8 ),        # Top
        Inches(4.2 ),        # Width
        Inches(0.4)         # Height 
    )

    # Borda azul escura (mais fina)
    legenda.line.color.rgb = RGBColor(0, 51, 102)
    legenda.line.width = Pt(0.9)

    tf_legenda = legenda.text_frame
    tf_legenda.clear()
    tf_legenda.margin_left = Inches(0.15)
    tf_legenda.margin_top = Inches(0.05)


    # -------- ITENS (HORIZONTAIS) -------- #
    legenda_itens = [
        ("Muito Alto",  RGBColor(192, 0, 0)),
        ("Alto",        RGBColor(237, 125, 49)),
        ("Médio",       RGBColor(255, 192, 0)),
        ("Baixo",       RGBColor(0, 97, 0)),
        ("Muito Baixo", RGBColor(146, 208, 80)),
    ]

    start_x = Inches(0.3)
    y = Inches(4.9)
    espacamento = Inches(0.8)

    for texto, cor in legenda_itens:
        # Círculo grande
        circulo = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            start_x, y,
            Inches(0.18), Inches(0.18)
        )
        circulo.fill.solid()
        circulo.fill.fore_color.rgb = cor
        circulo.line.fill.background()

        # Texto ao lado do círculo
        label = slide.shapes.add_textbox(
            start_x + Inches(0.12),
            y - Inches(0.02),
            Inches(0.6),
            Inches(0.25)
        )
        p = label.text_frame.paragraphs[0]
        p.text = texto
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.LEFT

        start_x += espacamento
      
    titulo_elementos = slide.shapes.add_textbox(
        Inches(4.7),     # à direita da legenda anterior
        Inches(4.55),
        Inches(4.8),
        Inches(0.25)
    )

    p = titulo_elementos.text_frame.paragraphs[0]
    p.text = "Elementos"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.LEFT


    # ========= CAIXA DA LEGENDA ========= #
    legenda_elementos = slide.shapes.add_textbox(
        Inches(4.5),     # Left
        Inches(4.8),     # Top
        Inches(4.6),     # Width
        Inches(0.4)      # Height
    )

    # Fundo azul escuro
    legenda_elementos.fill.solid()
    legenda_elementos.fill.fore_color.rgb = RGBColor(30, 115, 190)

    # Borda
    legenda_elementos.line.color.rgb = RGBColor(0, 51, 102)
    legenda_elementos.line.width = Pt(0.9)


    # ========= CAMINHO DAS IMAGENS ========= #
    current_dir = os.path.dirname(__file__)       # pptx-generator/slides/
    parent_dir = os.path.dirname(current_dir)     # pptx-generator/
    images_dir = os.path.join(parent_dir, "images")


    # ========= ITENS DA LEGENDA ========= #
    elementos = [
        ("ICON_TECNOLOGIA.png",   "Tecnologia"),
        ("ICON_PROCESSOS.png",    "Processos"),
        ("ICON_PESSOAS.png",      "Pessoas"),
        ("ICON_INFORMACOES.png",  "Informação"),
        ("ICON_GESTAO.png",       "Gestão"),
    ]

    start_x = Inches(4.6)
    y = Inches(4.9)
    espacamento = Inches(0.95)

    for nome_icone, texto in elementos:
        image_path = os.path.join(images_dir, nome_icone)

        # Ícone
        try:
            slide.shapes.add_picture(
                image_path,
                start_x,
                y,
                width=Inches(0.28)
            )
        except Exception as e:
            print(f"Erro ao carregar ícone {nome_icone}: {e}", file=sys.stderr)

        # Texto ao lado do ícone
        label = slide.shapes.add_textbox(
            start_x + Inches(0.18),
            y - Inches(0.02),
            Inches(0.8),
            Inches(0.3)
        )

        p = label.text_frame.paragraphs[0]
        p.text = texto
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.LEFT

        start_x += espacamento
    # === LOGO === #
    logo = dados.get("imagens", {}).get("logo_empresa")
    if logo:
        try:
            if ',' in logo:
                logo = logo.split(',')[1]
            stream = BytesIO(base64.b64decode(logo))
            slide.shapes.add_picture(stream, Inches(9.1), Inches(4.5), width=Inches(0.8))
        except Exception as e:
            print(f"Erro ao adicionar logo: {e}", file=sys.stderr)

    return slide

def gerar_vulnerabilidades(pres, dados):

    if not dados.get("vulnerabilidades"):
        itens = [
            {
                "elemento": f"Elemento {i}",
                "nc": f"NC-{i}",
                "nao_conformidade": f"Descrição da NC {i}",
                "observacao": f"Obs {i}",
                "risco": "Médio",
                "localizacao": "Área X",
                "criticidade": "Alta"
            }
            for i in range(1, 21)
        ]
    else:
        itens = dados["vulnerabilidades"]

    
    itens_por_slide = 5

    for i in range(0, len(itens), itens_por_slide):
        lote = itens[i:i + itens_por_slide]
        criar_slide_vulnerabilidades(pres, dados, lote)

    print("✅ Slides de Vulnerabilidades criados!", file=sys.stderr)
