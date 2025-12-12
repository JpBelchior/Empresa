#!/usr/bin/env python3
# -*- coding: utf-8 -*-
import sys
import os
import base64
from io import BytesIO
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


def gerar_seguranca(pres, dados):
    """Gera o slide de Segurança Pública"""
    print("Criando slide de Segurança Pública...", file=sys.stderr, flush=True)
    slide = pres.slides.add_slide(pres.slide_layouts[6])

    # === FUNDO BRANCO === #
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # === PARALELOGRAMO AZUL (número 4) === #
    para_azul = slide.shapes.add_shape(
        MSO_SHAPE.PARALLELOGRAM,
        Inches(0.2),
        Inches(0.45),
        Inches(0.5),
        Inches(0.5)
    )

    para_azul.fill.solid()
    para_azul.fill.fore_color.rgb = RGBColor(0, 102, 204)
    para_azul.line.fill.background()

    tf = para_azul.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "4"
    p.font.name = "Arial"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # === PARALELOGRAMO CINZA (SEGURANÇA PÚBLICA) === #
    para_cinza = slide.shapes.add_shape(
        MSO_SHAPE.PARALLELOGRAM,
        Inches(0.85),
        Inches(0.45),
        Inches(4.0),
        Inches(0.5)
    )

    para_cinza.fill.solid()
    para_cinza.fill.fore_color.rgb = RGBColor(64, 64, 64)
    para_cinza.line.fill.background()

    tf = para_cinza.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "SEGURANÇA PÚBLICA"
    p.font.name = "Arial"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.LEFT
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.2)

    # === IMAGEM SEGURANÇA PÚBLICA === #
    current_dir = os.path.dirname(__file__)
    parent_dir = os.path.dirname(current_dir)
    imagem_path = os.path.join(parent_dir, "images", "seguranca.png")

    try:
        if os.path.exists(imagem_path):
            slide.shapes.add_picture(
                imagem_path,
                Inches(0.3),      # Margem esquerda
                Inches(1.35),     # Abaixo do título
                width=Inches(8.7),
                height=Inches(3.8)
            )
        else:
            print(f"Imagem não encontrada: {imagem_path}", file=sys.stderr, flush=True)
    except Exception as e:
        print(f"Erro ao adicionar imagem de segurança: {e}", file=sys.stderr, flush=True)

    # === LOGO DA EMPRESA (canto inferior direito) === #
    logo = dados.get("imagens", {}).get("logo_empresa")
    if logo:
        try:
            if ',' in logo:
                logo = logo.split(',')[1]
            
            img_data = base64.b64decode(logo)
            stream = BytesIO(img_data)
            
            slide.shapes.add_picture(
                stream,
                Inches(9.0),   # Canto direito
                Inches(4.5),   # Canto inferior
                width=Inches(0.9)
            )
        except Exception as e:
            print(f"Erro ao adicionar logo: {e}", file=sys.stderr, flush=True)

    print("Slide de Segurança Pública criado!", file=sys.stderr, flush=True)
    return slide